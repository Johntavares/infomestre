from pptx import Presentation
import sys

def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        for i, slide in enumerate(prs.slides):
            print(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(shape.text)
            print()
    except Exception as e:
        print(f"Error reading pptx: {e}")

if __name__ == "__main__":
    if len(sys.argv) > 1:
        extract_text_from_pptx(sys.argv[1])
    else:
        print("Provide path to pptx")
